"""
Extract dominant colors from reference images and PowerPoint slides.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import Counter
import io


def extract_colors_from_image(image_path, num_colors=5):
    """Extract dominant colors from an image."""
    img = Image.open(image_path)
    img = img.convert('RGB')
    
    # Resize for faster processing
    img.thumbnail((150, 150))
    
    # Get all pixels
    pixels = list(img.getdata())
    
    # Count most common colors
    color_counts = Counter(pixels)
    most_common = color_counts.most_common(num_colors)
    
    return [color for color, count in most_common]


def extract_colors_from_pptx(pptx_path, slide_index=0, num_colors=5):
    """Extract colors from images in a PowerPoint slide."""
    prs = Presentation(pptx_path)
    
    if slide_index >= len(prs.slides):
        print(f"Slide {slide_index} not found in presentation.")
        return []
    
    slide = prs.slides[slide_index]
    colors = []
    
    # Extract images from slide
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            # This is an image shape
            image = shape.image
            image_bytes = image.blob
            img = Image.open(io.BytesIO(image_bytes))
            img = img.convert('RGB')
            img.thumbnail((150, 150))
            
            pixels = list(img.getdata())
            color_counts = Counter(pixels)
            most_common = color_counts.most_common(num_colors)
            colors.extend([color for color, count in most_common])
    
    return colors


def rgb_to_hex(rgb):
    """Convert RGB tuple to hex color string."""
    return '#{:02x}{:02x}{:02x}'.format(rgb[0], rgb[1], rgb[2])


def extract_all_reference_colors():
    """Extract colors from all reference files."""
    colors = []
    
    # Extract from Picture1.png
    if os.path.exists('Picture1.png'):
        print("Extracting colors from Picture1.png...")
        colors.extend(extract_colors_from_image('Picture1.png', 5))
    
    # Extract from Picture2.jpg
    if os.path.exists('Picture2.jpg'):
        print("Extracting colors from Picture2.jpg...")
        colors.extend(extract_colors_from_image('Picture2.jpg', 5))
    
    # Extract from ye.pptx first slide
    if os.path.exists('ye.pptx'):
        print("Extracting colors from ye.pptx...")
        colors.extend(extract_colors_from_pptx('ye.pptx', 0, 5))
    
    # Remove duplicates while preserving order
    unique_colors = []
    seen = set()
    for color in colors:
        if color not in seen:
            seen.add(color)
            unique_colors.append(color)
    
    return unique_colors


if __name__ == '__main__':
    colors = extract_all_reference_colors()
    print(f"\nExtracted {len(colors)} unique colors:")
    for i, color in enumerate(colors[:15]):  # Show first 15
        hex_color = rgb_to_hex(color)
        print(f"{i+1}. RGB{color} -> {hex_color}")
